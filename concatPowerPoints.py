import os
import comtypes.client
from spire.presentation import *
from spire.presentation.common import *
from tkinter import filedialog
from tkinter import messagebox
from tkinter import *
from pptx import Presentation as pptxPresentation
WATERMARK_TEXT = "Evaluation Warning : The document was created with Spire.Presentation for Python"

def create_one_presentation(dir, pptx_path):
    combined_present = Presentation()
    if combined_present.Slides.Count > 0:
        combined_present.Slides.RemoveAt(0)
    powerpoint_extensions = (".ppt", ".pptx", ".pps", ".ppsx", ".pot", ".potx")
    for file in os.listdir(dir):
        if file.lower().endswith(powerpoint_extensions):
            if file.startswith("~$"):
                continue
            filepath = os.path.join(dir, file)
            try:
                temp_present = Presentation()
                temp_present.LoadFromFile(filepath)

                for slide in temp_present.Slides:
                    combined_present.Slides.AppendBySlide(slide)
                
                temp_present.Dispose()
            except Exception as e:
                messagebox.showerror("Error Loading Presentation", f"❌ Failed to load '{file}':\n{e}")
    combined_present.SaveToFile(pptx_path, FileFormat.Pptx2016)
    combined_present.Dispose()
    print("✅ merged everything!")

def remove_watermark_text(pptx_path, output_path, target_text):
    prsent = pptxPresentation(pptx_path)

    for slide in prsent.slides:
        for shape in slide.shapes: # search for watermarks and the shape around it
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if target_text in run.text:
                        run.text = run.text.replace(target_text, " ")
                        sp = shape._element
                        sp.getparent().remove(sp)

    prsent.save(output_path)
    if os.path.exists(pptx_path):
        os.remove(pptx_path)

def pptx_to_pdf(input_path, output_path):
    print(f"🧪 Converting {input_path} to PDF...")
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    try:
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        presentation.SaveAs(output_path, 32)  # 32 = PDF
        presentation.Close()
        print(f"✅ Converted to PDF: {output_path}")
    except Exception as e:
        messagebox.showerror("PDF Conversion Failed", f"❌ Conversion failed:\n{e}")
    finally:
        powerpoint.Quit()

if __name__ == "__main__":
    root = Tk()
    root.withdraw()
    source_directory = filedialog.askdirectory()
    if not source_directory:
        messagebox.showwarning("No Folder Selected", "No folder was selected.")
    else:
        pptx_path = os.path.abspath(os.path.join(source_directory, "mergedPresentation.pptx"))
        cleaned_path = pptx_path.replace(".pptx", "_cleaned.pptx")
        pdf_path = pptx_path.replace(".pptx", ".pdf")
        create_one_presentation(source_directory, pptx_path)
        remove_watermark_text(pptx_path, cleaned_path, WATERMARK_TEXT)
        pptx_to_pdf(cleaned_path, pdf_path)